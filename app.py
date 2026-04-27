import argparse
import json
import os
import re
import shutil
from dataclasses import asdict, dataclass, field
from datetime import datetime
from heapq import nlargest
from pathlib import Path
from threading import Event, Lock, Thread
from typing import Any, Iterable, List, Optional
from uuid import uuid4

from dotenv import load_dotenv
from flask import Flask, Response, jsonify, render_template, request
from openai import OpenAI
from werkzeug.exceptions import HTTPException
from werkzeug.utils import secure_filename


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".py",
    ".js",
    ".ts",
    ".tsx",
    ".jsx",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".csv",
    ".html",
    ".css",
    ".java",
    ".go",
    ".rs",
    ".c",
    ".cpp",
    ".h",
    ".hpp",
    ".sql",
    ".log",
}

IGNORE_DIRS = {
    ".git",
    ".venv",
    "venv",
    "node_modules",
    "__pycache__",
    ".idea",
    ".vscode",
    "dist",
    "build",
    "chunks",
}

DEFAULT_SOURCE_DIR = "./data"
DEFAULT_INDEX_PATH = "./data/index.json"
DEFAULT_CHUNKS_DIR = "./data/chunks"
DEFAULT_KB_DIR = "./data/knowledge_bases"
DEFAULT_KB_META_PATH = "./data/knowledge_bases.json"
DEFAULT_EXCLUDED_FILE_NAMES = {
    "index.json",
    "knowledge_bases.json",
}
DEFAULT_CHUNK_SIZE = 1000
DEFAULT_CHUNK_OVERLAP = 200
DEFAULT_TOP_K = 5
DEFAULT_QWEN_MODEL = "qwen-plus"
DEFAULT_QWEN_BASE_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1"

DOCUMENT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".tsx", ".jsx", ".json",
    ".yaml", ".yml", ".toml", ".ini", ".csv", ".html", ".css",
    ".java", ".go", ".rs", ".c", ".cpp", ".h", ".hpp", ".sql", ".log",
}

OFFICE_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}

UPLOAD_FOLDER = "./data/uploads"
ALLOWED_EXTENSIONS = DOCUMENT_EXTENSIONS | OFFICE_EXTENSIONS
MAX_CONTENT_LENGTH = 50 * 1024 * 1024


def get_vectorizer_cache_path(index_path: str) -> str:
    """获取向量化器缓存文件路径"""
    base = Path(index_path)
    return str(base.parent / f"{base.stem}_vectorizer.joblib")


def get_matrix_cache_path(index_path: str) -> str:
    """获取矩阵缓存文件路径"""
    base = Path(index_path)
    return str(base.parent / f"{base.stem}_matrix.npz")


def joblib_dump(obj, file_path: str, compress: int = 3) -> None:
    """使用 joblib 保存对象"""
    import joblib
    joblib.dump(obj, file_path, compress=compress)


def save_npz(file_path: str, matrix, compressed: bool = True) -> None:
    """保存 numpy 矩阵为 npz 格式"""
    import numpy as np
    if matrix is None:
        np.savez_compressed(file_path) if compressed else np.savez(file_path)
        return
    if compressed:
        np.savez_compressed(file_path, data=matrix)
    else:
        np.savez(file_path, data=matrix)


@dataclass
class DocumentChunk:
    path: str
    chunk_id: int
    text: str


@dataclass
class KnowledgeBaseInfo:
    id: str
    name: str
    description: str = ""
    created_at: str = ""
    document_count: int = 0
    chunk_count: int = 0

    def to_dict(self) -> dict[str, Any]:
        return {
            "id": self.id,
            "name": self.name,
            "description": self.description,
            "created_at": self.created_at,
            "document_count": self.document_count,
            "chunk_count": self.chunk_count,
        }


@dataclass
class BuildJob:
    job_id: str
    status: str
    source: str
    index_path: str
    kb_id: str = "default"
    chunks_dir: str = DEFAULT_CHUNKS_DIR
    chunk_size: int = DEFAULT_CHUNK_SIZE
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP
    message: str = "等待开始"
    result: dict[str, Any] | None = None
    error: str | None = None

    def to_dict(self) -> dict[str, Any]:
        return {
            "job_id": self.job_id,
            "status": self.status,
            "source": self.source,
            "index": self.index_path,
            "kb_id": self.kb_id,
            "chunks_dir": self.chunks_dir,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "message": self.message,
            "result": self.result,
            "error": self.error,
        }


class LocalKnowledgeBase:
    def __init__(self, chunk_size: int = DEFAULT_CHUNK_SIZE, chunk_overlap: int = DEFAULT_CHUNK_OVERLAP) -> None:
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.documents: List[DocumentChunk] = []
        self.vectorizer = None
        self.matrix = None

    def build(self, root_dir: str) -> int:
        self.documents = list(self._load_documents(root_dir))
        if not self.documents:
            raise ValueError(f"在目录 {root_dir} 中没有找到可索引的文本文件。")
        return len(self.documents)

    def save(self, index_path: str) -> None:
        if not self.documents:
            raise ValueError("索引尚未构建，不能保存。")

        payload = {
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "documents": [asdict(doc) for doc in self.documents],
        }

        path = Path(index_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
        joblib_dump(self.vectorizer, get_vectorizer_cache_path(index_path), compress=3)
        save_npz(get_matrix_cache_path(index_path), self.matrix, compressed=True)

    @classmethod
    def load(cls, index_path: str) -> "LocalKnowledgeBase":
        path = Path(index_path)
        if not path.exists():
            raise FileNotFoundError(f"索引文件不存在: {index_path}")

        payload = json.loads(path.read_text(encoding="utf-8"))
        kb = cls(
            chunk_size=payload.get("chunk_size", DEFAULT_CHUNK_SIZE),
            chunk_overlap=payload.get("chunk_overlap", DEFAULT_CHUNK_OVERLAP),
        )
        kb.documents = [
            DocumentChunk(**item)
            for item in payload["documents"]
            if Path(item["path"]).name not in DEFAULT_EXCLUDED_FILE_NAMES
        ]
        return kb

    def search(self, query: str, top_k: int = DEFAULT_TOP_K) -> List[dict]:
        if not self.documents:
            raise ValueError("知识库未初始化，请先构建或加载索引。")
        query = query.strip()
        if not query:
            return []

        query_terms = self._extract_query_terms(query)
        scored_results = []
        for doc in self.documents:
            score = self._score_document(query, query_terms, doc.text)
            if score <= 0:
                continue
            scored_results.append(
                {
                    "path": doc.path,
                    "chunk_id": doc.chunk_id,
                    "score": score,
                    "text": doc.text,
                }
            )

        top_results = nlargest(top_k, scored_results, key=lambda item: item["score"])
        if not top_results:
            return []

        max_score = max(item["score"] for item in top_results)
        if max_score <= 0:
            return []

        for item in top_results:
            item["score"] = round(float(item["score"] / max_score), 4)
        return top_results

    def _extract_query_terms(self, query: str) -> list[str]:
        normalized_query = self._normalize_text(query)
        terms: list[str] = []

        ascii_terms = re.findall(r"[a-z0-9]{2,}", normalized_query)
        terms.extend(ascii_terms)

        chinese_chars = [char for char in query if "\u4e00" <= char <= "\u9fff"]
        for size in (4, 3, 2):
            if len(chinese_chars) < size:
                continue
            for index in range(len(chinese_chars) - size + 1):
                terms.append("".join(chinese_chars[index : index + size]))

        if not terms:
            terms.append(normalized_query)

        deduped_terms: list[str] = []
        seen: set[str] = set()
        for term in terms:
            cleaned = term.strip()
            if not cleaned or cleaned in seen:
                continue
            seen.add(cleaned)
            deduped_terms.append(cleaned)
        return deduped_terms

    def _score_document(self, query: str, query_terms: list[str], text: str) -> float:
        normalized_text = self._normalize_text(text)
        normalized_query = self._normalize_text(query)
        score = 0.0

        if normalized_query and normalized_query in normalized_text:
            score += min(len(normalized_query), 24) * 3.0

        for term in query_terms:
            occurrences = normalized_text.count(term)
            if occurrences <= 0:
                continue
            score += occurrences * max(len(term), 1)

        return score

    def _normalize_text(self, text: str) -> str:
        return re.sub(r"\s+", " ", text.lower()).strip()

    def _load_documents(self, root_dir: str) -> Iterable[DocumentChunk]:
        root = Path(root_dir)
        if not root.exists():
            raise FileNotFoundError(f"目录不存在: {root_dir}")

        all_extensions = DOCUMENT_EXTENSIONS | OFFICE_EXTENSIONS

        for file_path in root.rglob("*"):
            if not file_path.is_file():
                continue
            if any(part in IGNORE_DIRS for part in file_path.parts):
                continue
            if file_path.name in DEFAULT_EXCLUDED_FILE_NAMES:
                continue
            if file_path.suffix.lower() not in all_extensions:
                continue

            text = self._read_text_file(file_path)
            if not text.strip():
                continue

            relative_path = str(file_path.relative_to(root))
            for chunk_id, chunk_text in enumerate(self._split_text(text)):
                yield DocumentChunk(path=relative_path, chunk_id=chunk_id, text=chunk_text)

    def _read_text_file(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()

        if suffix == ".pdf":
            return self._read_pdf_file(file_path)
        elif suffix in {".docx", ".doc"}:
            return self._read_docx_file(file_path)
        elif suffix in {".pptx", ".ppt"}:
            return self._read_pptx_file(file_path)
        elif suffix in {".xlsx", ".xls"}:
            return self._read_xlsx_file(file_path)

        encodings = ["utf-8", "utf-8-sig", "gbk", "gb18030"]
        for encoding in encodings:
            try:
                return file_path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
            except OSError:
                return ""
        return ""

    def _read_pdf_file(self, file_path: Path) -> str:
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() or ""
                return text
        except ImportError:
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    text = ""
                    for page in pdf.pages:
                        text += page.extract_text() or ""
                    return text
            except ImportError:
                return ""
        except Exception:
            return ""

    def _read_docx_file(self, file_path: Path) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except ImportError:
            return ""
        except Exception:
            return ""

    def _read_pptx_file(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except ImportError:
            return ""
        except Exception:
            return ""

    def _read_xlsx_file(self, file_path: Path) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
            return text
        except ImportError:
            return ""
        except Exception:
            return ""

    def _split_text(self, text: str) -> List[str]:
        text = text.replace("\r\n", "\n").strip()
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        start = 0
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(text):
                break
            start = max(end - self.chunk_overlap, start + 1)
        return chunks

    def save_chunks_to_files(self, root_dir: str, chunks_dir: str) -> None:
        chunks_path = Path(chunks_dir)
        chunks_path.mkdir(parents=True, exist_ok=True)

        root = Path(root_dir)
        file_chunk_counts = {}

        for doc in self.documents:
            source_file = doc.path
            if source_file not in file_chunk_counts:
                file_chunk_counts[source_file] = 0

            safe_filename = source_file.replace(os.sep, "_").replace(" ", "_")
            chunk_file = chunks_path / f"{safe_filename}_chunk{doc.chunk_id}.txt"
            chunk_file.write_text(doc.text, encoding="utf-8")
            file_chunk_counts[source_file] += 1


class QwenChat:
    def __init__(self, api_key: str, model: str, base_url: str) -> None:
        self.client = OpenAI(api_key=api_key, base_url=base_url)
        self.model = model

    def answer(self, question: str, contexts: List[dict]) -> str:
        context_text = "\n\n".join(
            [
                f"[文件: {item['path']} | 分片: {item['chunk_id']} | 相似度: {item['score']:.4f}]\n{item['text']}"
                for item in contexts
            ]
        )

        messages = [
            {
                "role": "system",
                "content": (
                    "你是一个本地知识库问答助手。"
                    "只能优先依据提供的检索上下文回答。"
                    "如果上下文不足，请明确说明‘根据当前知识库内容无法确定’，不要编造事实。"
                ),
            },
            {
                "role": "user",
                "content": (
                    f"用户问题：{question}\n\n"
                    f"检索上下文：\n{context_text}\n\n"
                    "请基于以上内容给出中文回答，并在结尾列出引用到的文件路径。"
                ),
            },
        ]

        response = self.client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=0.2,
        )
        return response.choices[0].message.content or "模型没有返回内容。"


class KnowledgeBaseService:
    def __init__(self) -> None:
        self._lock = Lock()
        self.kb: LocalKnowledgeBase | None = None
        self.loaded_index_path: str | None = None
        self.loaded_kb_id: str | None = None
        self.build_jobs: dict[str, BuildJob] = {}
        self.active_build_job_id: str | None = None
        self.loading_indexes: dict[str, Event] = {}
        self.loading_errors: dict[str, str] = {}
        self.warmup_started = False
        self.warmup_error: str | None = None
        self._ensure_default_kb()

    def _ensure_default_kb(self) -> None:
        Path(DEFAULT_KB_DIR).mkdir(parents=True, exist_ok=True)
        Path(UPLOAD_FOLDER).mkdir(parents=True, exist_ok=True)
        
        meta_path = Path(DEFAULT_KB_META_PATH)
        if not meta_path.exists():
            default_kb = KnowledgeBaseInfo(
                id="default",
                name="剑来知识库",
                description="剑来小说1-500章本地知识库",
                created_at=datetime.now().isoformat(),
            )
            self._save_kb_meta([default_kb])

    def _load_kb_meta(self) -> List[KnowledgeBaseInfo]:
        meta_path = Path(DEFAULT_KB_META_PATH)
        if not meta_path.exists():
            return []
        try:
            data = json.loads(meta_path.read_text(encoding="utf-8"))
            return [KnowledgeBaseInfo(**item) for item in data.get("bases", [])]
        except Exception:
            return []

    def _save_kb_meta(self, bases: List[KnowledgeBaseInfo]) -> None:
        meta_path = Path(DEFAULT_KB_META_PATH)
        meta_path.parent.mkdir(parents=True, exist_ok=True)
        data = {"bases": [kb.to_dict() for kb in bases]}
        meta_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

    def _get_kb_path(self, kb_id: str) -> str:
        return str(Path(DEFAULT_KB_DIR) / kb_id)

    def _get_kb_index_path(self, kb_id: str) -> str:
        return str(Path(DEFAULT_KB_DIR) / kb_id / "index.json")

    def _get_kb_chunks_dir(self, kb_id: str) -> str:
        return str(Path(DEFAULT_KB_DIR) / kb_id / "chunks")

    def _get_kb_upload_dir(self, kb_id: str) -> str:
        return str(Path(UPLOAD_FOLDER) / kb_id)

    def list_knowledge_bases(self) -> List[dict]:
        bases = self._load_kb_meta()
        result = []
        for kb in bases:
            index_path = Path(self._get_kb_index_path(kb.id))
            kb_info = kb.to_dict()
            kb_info["index_exists"] = index_path.exists()
            result.append(kb_info)
        return result

    def create_knowledge_base(self, name: str, description: str = "") -> dict:
        kb_id = uuid4().hex[:8]
        kb = KnowledgeBaseInfo(
            id=kb_id,
            name=name,
            description=description,
            created_at=datetime.now().isoformat(),
        )
        bases = self._load_kb_meta()
        bases.append(kb)
        self._save_kb_meta(bases)
        
        kb_dir = Path(self._get_kb_path(kb_id))
        kb_dir.mkdir(parents=True, exist_ok=True)
        (kb_dir / "chunks").mkdir(parents=True, exist_ok=True)
        
        upload_dir = Path(self._get_kb_upload_dir(kb_id))
        upload_dir.mkdir(parents=True, exist_ok=True)
        
        return kb.to_dict()

    def delete_knowledge_base(self, kb_id: str) -> bool:
        if kb_id == "default":
            raise ValueError("不能删除默认知识库")
        
        bases = self._load_kb_meta()
        bases = [kb for kb in bases if kb.id != kb_id]
        self._save_kb_meta(bases)
        
        kb_dir = Path(self._get_kb_path(kb_id))
        if kb_dir.exists():
            shutil.rmtree(kb_dir)
        
        upload_dir = Path(self._get_kb_upload_dir(kb_id))
        if upload_dir.exists():
            shutil.rmtree(upload_dir)
        
        return True

    def upload_file(self, kb_id: str, file) -> dict:
        upload_dir = Path(self._get_kb_upload_dir(kb_id))
        upload_dir.mkdir(parents=True, exist_ok=True)
        
        filename = secure_filename(file.filename)
        file_path = upload_dir / filename
        file.save(str(file_path))
        
        return {
            "filename": filename,
            "path": str(file_path),
            "size": file_path.stat().st_size,
        }

    def list_uploaded_files(self, kb_id: str) -> List[dict]:
        if kb_id == "default":
            source_dir = Path(DEFAULT_SOURCE_DIR)
        else:
            source_dir = Path(self._get_kb_upload_dir(kb_id))
        
        if not source_dir.exists():
            return []
        
        all_extensions = DOCUMENT_EXTENSIONS | OFFICE_EXTENSIONS
        files = []
        for f in source_dir.rglob("*"):
            if not f.is_file():
                continue
            if f.name in DEFAULT_EXCLUDED_FILE_NAMES:
                continue
            if any(part in IGNORE_DIRS for part in f.parts):
                continue
            if f.suffix.lower() not in all_extensions:
                continue
            files.append({
                "filename": f.name,
                "path": str(f),
                "size": f.stat().st_size,
            })
        return files

    def delete_uploaded_file(self, kb_id: str, filename: str) -> bool:
        file_path = Path(self._get_kb_upload_dir(kb_id)) / secure_filename(filename)
        if file_path.exists() and file_path.is_file():
            file_path.unlink()
            return True
        return False

    def build_index(
        self,
        source: str,
        index_path: str,
        kb_id: str = "default",
        chunks_dir: str = DEFAULT_CHUNKS_DIR,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
    ) -> dict:
        kb = LocalKnowledgeBase(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        count = kb.build(source)
        kb.save(index_path)
        kb.save_chunks_to_files(source, chunks_dir)
        
        bases = self._load_kb_meta()
        for base in bases:
            if base.id == kb_id:
                base.chunk_count = count
                if kb_id == "default":
                    base.document_count = len(self.list_uploaded_files(kb_id))
                else:
                    upload_dir = Path(self._get_kb_upload_dir(kb_id))
                    base.document_count = len(list(upload_dir.glob("*"))) if upload_dir.exists() else 0
                break
        self._save_kb_meta(bases)
        
        with self._lock:
            self.kb = kb
            self.loaded_index_path = index_path
            self.loaded_kb_id = kb_id
        return {
            "source": source,
            "index": index_path,
            "kb_id": kb_id,
            "chunks_dir": chunks_dir,
            "chunks": count,
            "chunk_size": chunk_size,
            "chunk_overlap": chunk_overlap,
        }

    def start_build_index(
        self,
        source: str,
        index_path: str,
        kb_id: str = "default",
        chunks_dir: str = DEFAULT_CHUNKS_DIR,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
    ) -> dict[str, Any]:
        with self._lock:
            active_job = self.build_jobs.get(self.active_build_job_id or "")
            if active_job and active_job.status in {"queued", "running"}:
                same_request = (
                    active_job.source == source
                    and active_job.index_path == index_path
                    and active_job.chunk_size == chunk_size
                    and active_job.chunk_overlap == chunk_overlap
                )
                if same_request:
                    return active_job.to_dict()
                raise ValueError("已有索引构建任务进行中，请等待当前任务完成后再试。")

            job = BuildJob(
                job_id=uuid4().hex,
                status="queued",
                source=source,
                index_path=index_path,
                kb_id=kb_id,
                chunks_dir=chunks_dir,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                message="任务已创建，等待开始构建。",
            )
            self.build_jobs[job.job_id] = job
            self.active_build_job_id = job.job_id

        Thread(target=self._run_build_job, args=(job.job_id,), daemon=True).start()
        return job.to_dict()

    def _run_build_job(self, job_id: str) -> None:
        with self._lock:
            job = self.build_jobs[job_id]
            job.status = "running"
            job.message = "正在扫描文件并构建索引，请稍候..."

        try:
            result = self.build_index(
                source=job.source,
                index_path=job.index_path,
                kb_id=job.kb_id,
                chunks_dir=job.chunks_dir,
                chunk_size=job.chunk_size,
                chunk_overlap=job.chunk_overlap,
            )
            with self._lock:
                current_job = self.build_jobs[job_id]
                current_job.status = "completed"
                current_job.message = f"索引构建完成：共 {result['chunks']} 个分片。"
                current_job.result = result
                current_job.error = None
        except Exception as error:
            with self._lock:
                current_job = self.build_jobs[job_id]
                current_job.status = "failed"
                current_job.message = f"构建失败：{error}"
                current_job.error = str(error)
                current_job.result = None
        finally:
            with self._lock:
                if self.active_build_job_id == job_id:
                    self.active_build_job_id = None

    def get_build_job(self, job_id: str) -> dict[str, Any]:
        with self._lock:
            job = self.build_jobs.get(job_id)
            if job is None:
                raise FileNotFoundError(f"构建任务不存在: {job_id}")
            return job.to_dict()

    def ensure_loaded(self, index_path: str) -> LocalKnowledgeBase:
        wait_event: Event | None = None
        should_load = False
        with self._lock:
            if self.kb is not None and self.loaded_index_path == index_path:
                return self.kb

            wait_event = self.loading_indexes.get(index_path)
            if wait_event is None:
                wait_event = Event()
                self.loading_indexes[index_path] = wait_event
                self.loading_errors.pop(index_path, None)
                should_load = True

        if not should_load:
            wait_event.wait()
            with self._lock:
                if self.kb is not None and self.loaded_index_path == index_path:
                    return self.kb
                error_message = self.loading_errors.get(index_path)
            if error_message:
                raise RuntimeError(error_message)
            raise FileNotFoundError(f"索引文件不存在: {index_path}")

        try:
            kb = LocalKnowledgeBase.load(index_path)
            with self._lock:
                self.kb = kb
                self.loaded_index_path = index_path
            return kb
        except Exception as error:
            with self._lock:
                self.loading_errors[index_path] = str(error)
            raise
        finally:
            with self._lock:
                current_event = self.loading_indexes.pop(index_path, None)
                if current_event is not None:
                    current_event.set()

    def start_warmup(self, kb_id: str = "default") -> None:
        with self._lock:
            if self.warmup_started:
                return
            self.warmup_started = True

        Thread(target=self._warmup_index, args=(kb_id,), daemon=True).start()

    def _warmup_index(self, kb_id: str) -> None:
        try:
            index_path = self._get_kb_index_path(kb_id)
            
            if kb_id == "default":
                source_dir = DEFAULT_SOURCE_DIR
            else:
                source_dir = self._get_kb_upload_dir(kb_id)
            
            if Path(index_path).exists():
                self.ensure_loaded(index_path)
                with self._lock:
                    self.loaded_kb_id = kb_id
                print(f"[预热] 已加载现有索引: {index_path}")
            elif Path(source_dir).exists():
                all_extensions = DOCUMENT_EXTENSIONS | OFFICE_EXTENSIONS
                has_files = any(
                    f.suffix.lower() in all_extensions
                    for f in Path(source_dir).rglob("*")
                    if f.is_file() and f.name not in DEFAULT_EXCLUDED_FILE_NAMES
                    and not any(part in IGNORE_DIRS for part in f.parts)
                )
                if has_files:
                    print(f"[预热] 索引文件不存在，开始自动构建索引...")
                    print(f"[预热] 扫描目录: {source_dir}")
                    result = self.build_index(
                        source=source_dir,
                        index_path=index_path,
                        kb_id=kb_id,
                        chunks_dir=self._get_kb_chunks_dir(kb_id),
                    )
                    print(f"[预热] 索引构建完成：共 {result['chunks']} 个分片")
                else:
                    print(f"[预热] 知识库 {kb_id} 暂无文件，跳过构建")
            else:
                print(f"[预热] 知识库 {kb_id} 目录不存在，跳过构建")
        except Exception as error:
            with self._lock:
                self.warmup_error = str(error)
            print(f"[预热] 错误: {error}")

    def ask(self, question: str, kb_id: str = "default", top_k: int = DEFAULT_TOP_K) -> dict:
        load_dotenv()
        index_path = self._get_kb_index_path(kb_id)
        kb = self.ensure_loaded(index_path)
        with self._lock:
            self.loaded_kb_id = kb_id
        results = kb.search(question, top_k=top_k)
        if not results:
            return {
                "answer": "未检索到相关内容，请尝试更换问题或重新建立索引。",
                "results": [],
            }

        api_key = get_env("DASHSCOPE_API_KEY")
        model = get_env("QWEN_MODEL", DEFAULT_QWEN_MODEL)
        base_url = get_env("QWEN_BASE_URL", DEFAULT_QWEN_BASE_URL)

        chat = QwenChat(api_key=api_key, model=model, base_url=base_url)
        answer = chat.answer(question, results)
        return {
            "answer": answer,
            "results": results,
        }

    def status(self, kb_id: str = "default") -> dict:
        index_path = self._get_kb_index_path(kb_id)
        index_file = Path(index_path)
        index_exists = index_file.exists()
        loaded = self.kb is not None and self.loaded_kb_id == kb_id
        document_count = len(self.kb.documents) if loaded and self.kb is not None else 0
        
        bases = self._load_kb_meta()
        kb_info = None
        for kb in bases:
            if kb.id == kb_id:
                kb_info = kb.to_dict()
                break
        
        active_job = self.build_jobs.get(self.active_build_job_id or "")
        warming_up = index_path in self.loading_indexes
        return {
            "kb_id": kb_id,
            "index_exists": index_exists,
            "index_path": index_path,
            "loaded": loaded,
            "document_count": document_count,
            "api_key_configured": bool(os.getenv("DASHSCOPE_API_KEY")),
            "build_job": active_job.to_dict() if active_job else None,
            "warming_up": warming_up,
            "warmup_error": self.warmup_error if kb_id == "default" else None,
            "kb_info": kb_info,
        }


def get_env(name: str, default: str | None = None) -> str:
    value = os.getenv(name, default)
    if not value:
        raise ValueError(f"缺少环境变量: {name}")
    return value


service = KnowledgeBaseService()
service.start_warmup()
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

port = int(os.getenv('PORT', 7860))


@app.get("/")
def index() -> str:
    return render_template("index.html")


@app.get("/favicon.ico")
def favicon() -> Response:
    return Response(status=204)


@app.get("/.well-known/appspecific/com.chrome.devtools.json")
def chrome_devtools_manifest() -> Response:
    return Response(status=204)


@app.get("/api/kb/list")
def api_kb_list():
    return jsonify(service.list_knowledge_bases())


@app.post("/api/kb/create")
def api_kb_create():
    payload = request.get_json(silent=True) or {}
    name = (payload.get("name") or "").strip()
    description = (payload.get("description") or "").strip()
    
    if not name:
        return jsonify({"error": "知识库名称不能为空"}), 400
    
    result = service.create_knowledge_base(name=name, description=description)
    return jsonify(result), 201


@app.delete("/api/kb/<kb_id>")
def api_kb_delete(kb_id: str):
    try:
        service.delete_knowledge_base(kb_id)
        return jsonify({"success": True})
    except ValueError as e:
        return jsonify({"error": str(e)}), 400


@app.get("/api/kb/<kb_id>/files")
def api_kb_files(kb_id: str):
    files = service.list_uploaded_files(kb_id)
    return jsonify(files)


@app.post("/api/kb/<kb_id>/upload")
def api_kb_upload(kb_id: str):
    if 'file' not in request.files:
        return jsonify({"error": "没有上传文件"}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "没有选择文件"}), 400
    
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": f"不支持的文件格式: {ext}"}), 400
    
    result = service.upload_file(kb_id, file)
    return jsonify(result), 201


@app.delete("/api/kb/<kb_id>/files/<filename>")
def api_kb_file_delete(kb_id: str, filename: str):
    success = service.delete_uploaded_file(kb_id, filename)
    if success:
        return jsonify({"success": True})
    return jsonify({"error": "文件不存在"}), 404


@app.get("/api/status")
def api_status():
    load_dotenv()
    kb_id = request.args.get("kb_id", "default")
    return jsonify(service.status(kb_id=kb_id))


@app.post("/api/build")
def api_build():
    load_dotenv()
    payload = request.get_json(silent=True) or {}
    kb_id = payload.get("kb_id", "default")
    chunk_size = int(payload.get("chunk_size", DEFAULT_CHUNK_SIZE))
    chunk_overlap = int(payload.get("chunk_overlap", DEFAULT_CHUNK_OVERLAP))

    upload_dir = service._get_kb_upload_dir(kb_id)
    index_path = service._get_kb_index_path(kb_id)
    chunks_dir = service._get_kb_chunks_dir(kb_id)

    result = service.start_build_index(
        source=upload_dir,
        index_path=index_path,
        kb_id=kb_id,
        chunks_dir=chunks_dir,
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    return jsonify(result), 202


@app.get("/api/build-status")
def api_build_status():
    job_id = (request.args.get("job_id") or "").strip()
    if not job_id:
        return jsonify({"error": "缺少 job_id 参数"}), 400
    return jsonify(service.get_build_job(job_id))


@app.post("/api/ask")
def api_ask():
    load_dotenv()
    payload = request.get_json(silent=True) or {}
    question = (payload.get("question") or "").strip()
    kb_id = payload.get("kb_id", "default")
    top_k = int(payload.get("top_k", DEFAULT_TOP_K))

    if not question:
        return jsonify({"error": "问题不能为空"}), 400

    result = service.ask(question=question, kb_id=kb_id, top_k=top_k)
    return jsonify(result)


@app.errorhandler(Exception)
def handle_exception(error):
    if isinstance(error, HTTPException):
        response = error.get_response()
        response.data = json.dumps(
            {
                "error": error.description,
                "code": error.code,
                "name": error.name,
            },
            ensure_ascii=False,
        )
        response.content_type = "application/json; charset=utf-8"
        return response

    return jsonify({"error": str(error)}), 500


def build_command(args: argparse.Namespace) -> None:
    result = service.build_index(
        source=args.source,
        index_path=args.index,
        chunk_size=args.chunk_size,
        chunk_overlap=args.chunk_overlap,
    )
    print(
        f"已建立索引，共 {result['chunks']} 个文本分片，保存到 {result['index']}"
    )


def ask_command(args: argparse.Namespace) -> None:
    result = service.ask(question=args.question, index_path=args.index, top_k=args.top_k)
    print("=" * 80)
    print("检索结果：")
    for item in result["results"]:
        print(f"- {item['path']}#chunk-{item['chunk_id']} | score={item['score']:.4f}")
    print("=" * 80)
    print("回答：")
    print(result["answer"])


def serve_command(args: argparse.Namespace) -> None:
    load_dotenv()
    app.run(host=args.host, port=args.port, debug=False)


def create_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="本地知识库 + 千问远程模型问答")
    subparsers = parser.add_subparsers(dest="command", required=True)

    build_parser = subparsers.add_parser("build", help="扫描目录并建立本地索引")
    build_parser.add_argument("--source", default=DEFAULT_SOURCE_DIR, help="需要建立知识库的目录")
    build_parser.add_argument("--index", default=DEFAULT_INDEX_PATH, help="索引输出文件")
    build_parser.add_argument("--chunk-size", type=int, default=DEFAULT_CHUNK_SIZE, help="分片长度")
    build_parser.add_argument("--chunk-overlap", type=int, default=DEFAULT_CHUNK_OVERLAP, help="分片重叠长度")
    build_parser.set_defaults(func=build_command)

    ask_parser = subparsers.add_parser("ask", help="基于本地索引进行问答")
    ask_parser.add_argument("question", help="用户问题")
    ask_parser.add_argument("--index", default=DEFAULT_INDEX_PATH, help="已有索引文件")
    ask_parser.add_argument("--top-k", type=int, default=DEFAULT_TOP_K, help="检索返回分片数")
    ask_parser.set_defaults(func=ask_command)

    serve_parser = subparsers.add_parser("serve", help="启动 Web 服务")
    serve_parser.add_argument("--host", default="127.0.0.1", help="监听地址")
    serve_parser.add_argument("--port", type=int, default=7860, help="监听端口")
    serve_parser.set_defaults(func=serve_command)

    return parser


def main() -> None:
    load_dotenv()
    parser = create_parser()
    args = parser.parse_args()
    args.func(args)


if __name__ == "__main__":
    main()
